"""
One-off: parse e:\\1\\new_bot into docs/import-new-bot/
- 清扫机器人: PDF -> MD + embedded images extracted
- 酒店配送机器人: PPTX -> MD + images from zip
Also copies loose PNGs from cleaning folder into images/.
"""
from __future__ import annotations

import io
import re
import shutil
import zipfile
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation

SRC = Path(r"e:\1\new_bot")
OUT = Path(__file__).resolve().parent.parent / "docs" / "import-new-bot"


def slug(s: str) -> str:
    s = re.sub(r"[^\w\u4e00-\u9fff]+", "-", s.strip())
    return s.strip("-") or "unnamed"


def pdf_to_md_and_images(pdf_path: Path, images_dir: Path, md_parts: list[str]) -> None:
    images_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(pdf_path)
    title = pdf_path.stem
    md_parts.append(f"## PDF: {title}\n\n")
    for pi, page in enumerate(doc):
        text = page.get_text("text") or ""
        if text.strip():
            md_parts.append(f"### Page {pi + 1}\n\n{text.strip()}\n\n")
        for img_i, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            try:
                pix = fitz.Pixmap(doc, xref)
                if pix.n - pix.alpha > 3:
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                name = f"{slug(pdf_path.stem)}-p{pi + 1}-img{img_i + 1}.png"
                out_p = images_dir / name
                pix.save(out_p.as_posix())
                md_parts.append(f"![extracted](images/{name})\n\n")
            except Exception as e:  # noqa: BLE001
                md_parts.append(f"<!-- image extract failed p{pi + 1} img{img_i + 1}: {e} -->\n\n")
    doc.close()


def pptx_to_md_and_images(pptx_path: Path, images_dir: Path, md_parts: list[str]) -> None:
    images_dir.mkdir(parents=True, exist_ok=True)
    title = pptx_path.stem
    md_parts.append(f"## PPTX: {title}\n\n")

    # Text from slides
    prs = Presentation(pptx_path.as_posix())
    for si, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                lines.append(shape.text.strip())
        if lines:
            md_parts.append(f"### Slide {si}\n\n" + "\n\n".join(lines) + "\n\n")

    # Images from OOXML zip
    with zipfile.ZipFile(pptx_path, "r") as zf:
        media = [n for n in zf.namelist() if n.startswith("ppt/media/")]
        for n in sorted(media):
            ext = Path(n).suffix.lower() or ".bin"
            safe = slug(pptx_path.stem) + "-" + slug(Path(n).stem) + ext
            data = zf.read(n)
            (images_dir / safe).write_bytes(data)
            md_parts.append(f"![media](images/{safe})\n\n")


def main() -> None:
    if not SRC.is_dir():
        raise SystemExit(f"Missing source: {SRC}")

    children = sorted([p for p in SRC.iterdir() if p.is_dir()])
    if len(children) != 2:
        raise SystemExit(f"Expected 2 subfolders, got: {[p.name for p in children]}")

    # Folder 1: 清扫机器人
    clean_dir = children[0]
    deliv_dir = children[1]
    # Ensure order: 清扫 before 酒店 — sorted Unicode: 清 < 酒? Actually 清 is U+6E05, 酒 U+9152 — 清 < 酒 so cleaning first. Good.

    OUT.mkdir(parents=True, exist_ok=True)

    # --- Cleaning ---
    c_out = OUT / "01-cleaning-robot"
    c_img = c_out / "images"
    c_img.mkdir(parents=True, exist_ok=True)
    md1: list[str] = [
        "# 清扫机器人（新产品资料导入）\n\n",
        f"_来源目录：`{clean_dir.name}`_\n\n",
        "---\n\n",
    ]
    files = list(clean_dir.iterdir())
    pdfs = sorted([f for f in files if f.suffix.lower() == ".pdf"], key=lambda p: p.name.lower())
    imgs = sorted(
        [f for f in files if f.suffix.lower() in (".png", ".jpg", ".jpeg", ".webp")],
        key=lambda p: p.name.lower(),
    )
    vids = sorted([f for f in files if f.suffix.lower() in (".mp4", ".mov", ".MOV")], key=lambda p: p.name.lower())

    for f in pdfs:
        pdf_to_md_and_images(f, c_img, md1)
    if imgs:
        md1.append("## 目录内独立图片\n\n")
    for f in imgs:
        dest = c_img / f.name
        shutil.copy2(f, dest)
        md1.append(f"![source asset](images/{f.name})\n\n")
    if vids:
        md1.append("## 视频（未纳入仓库）\n\n")
    for f in vids:
        md1.append(f"<!-- 视频文件（未复制）: `{f.name}` -->\n\n")
    (c_out / "content.md").write_text("".join(md1), encoding="utf-8")

    # --- Hotel delivery ---
    d_out = OUT / "02-hotel-delivery-robot"
    d_img = d_out / "images"
    d_img.mkdir(parents=True, exist_ok=True)
    md2: list[str] = [
        "# 酒店配送机器人（新产品资料导入）\n\n",
        f"_来源目录：`{deliv_dir.name}`_\n\n",
        "---\n\n",
    ]
    d_files = list(deliv_dir.iterdir())
    d_pptx = sorted([f for f in d_files if f.suffix.lower() == ".pptx"], key=lambda p: p.name.lower())
    d_vids = sorted([f for f in d_files if f.suffix.lower() in (".mp4", ".mov", ".MOV")], key=lambda p: p.name.lower())
    for f in d_pptx:
        pptx_to_md_and_images(f, d_img, md2)
    if d_vids:
        md2.append("## 视频（未纳入仓库）\n\n")
    for f in d_vids:
        md2.append(f"<!-- 视频文件（未复制）: `{f.name}` -->\n\n")
    (d_out / "content.md").write_text("".join(md2), encoding="utf-8")

    # README index
    index = OUT / "README.md"
    index.write_text(
        "# new_bot 资料导入说明\n\n"
        "本目录由 `scripts/import-new-bot.py` 从 `e:\\\\1\\\\new_bot` 生成。\n\n"
        "| 子目录 | 对应产品方向 |\n"
        "|--------|----------------|\n"
        "| `01-cleaning-robot/` | 清扫机器人（PDF 正文 + PDF 内嵌图 + 目录内 PNG） |\n"
        "| `02-hotel-delivery-robot/` | 酒店配送机器人（PPTX 文本 + PPTX 内嵌媒体） |\n\n"
        "视频（`.mp4` / `.MOV`）未复制进仓库，仅在 `content.md` 中以注释列出文件名。\n",
        encoding="utf-8",
    )
    print("OK ->", OUT)


if __name__ == "__main__":
    main()
